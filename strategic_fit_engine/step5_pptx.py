"""
Step 5 — PowerPoint Summary Deck

Generates a concise executive-style .pptx from the scored targets and buyer profile.

Slides:
  1. Title slide
  2. Buyer profile & scoring criteria
  3. All targets ranked table
  4–6. One slide per top-3 company (snapshot + scores)
  7. DCF valuation summary
  8. Recommendation & next steps

Run standalone:
    python strategic_fit_engine/step5_pptx.py
"""
from __future__ import annotations

import json
import os
import re
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

DATA_DIR   = Path(__file__).parent / "data"
OUTPUT_DIR = Path(__file__).parent / "output"

# ── Brand colours ──────────────────────────────────────────────────────────────
NAVY  = RGBColor(0x25, 0x28, 0x50)
RED   = RGBColor(0xCC, 0x06, 0x05)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
GREY  = RGBColor(0x71, 0x80, 0x96)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
AMBER = RGBColor(0xF9, 0xA8, 0x25)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Low-level helpers ──────────────────────────────────────────────────────────

# XML 1.0 forbids control characters other than tab (0x09), LF (0x0A), CR (0x0D)
_XML_INVALID = re.compile(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]')

def _safe(text) -> str:
    """Strip XML-invalid characters so PowerPoint never shows a repair warning."""
    return _XML_INVALID.sub('', str(text) if not isinstance(text, str) else text)


def _rgb(r: int, g: int, b: int) -> RGBColor:
    return RGBColor(r, g, b)


def _solid(shape, color: RGBColor) -> None:
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def _strip_style(shape) -> None:
    """Remove the auto-added <p:style> element from a shape.

    python-pptx injects a <p:style> block that references theme colour slots
    (accent1, lt1, etc.). In a programmatically-built presentation those
    references are unresolved, which causes PowerPoint to show a repair
    dialog on every open. Removing the element eliminates the warning.
    """
    from pptx.oxml.ns import qn
    sp = shape._element
    style = sp.find(qn('p:style'))
    if style is not None:
        sp.remove(style)


def _add_rect(slide, l, t, w, h, color: RGBColor):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_AUTO_SHAPE_TYPE.RECTANGLE = 1
    _solid(shape, color)
    shape.line.fill.background()
    _strip_style(shape)
    return shape


def _tf(shape):
    return shape.text_frame


def _add_textbox(slide, text: str, l, t, w, h,
                 bold=False, size=14, color=None, align=PP_ALIGN.LEFT,
                 wrap=True) -> Any:
    from pptx.util import Pt
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = _safe(text)
    run.font.bold  = bold
    run.font.size  = Pt(size)
    run.font.color.rgb = color or NAVY
    return txb


def _heading(slide, text: str, top=Inches(0.28)):
    """Red left bar + white heading text on navy band."""
    _add_rect(slide, 0, top, SLIDE_W, Inches(0.52), NAVY)
    _add_rect(slide, 0, top, Inches(0.08), Inches(0.52), RED)
    _add_textbox(slide, text, Inches(0.22), top + Inches(0.06),
                 SLIDE_W - Inches(0.5), Inches(0.4),
                 bold=True, size=16, color=WHITE)


def _score_color(score: int, max_score: int) -> RGBColor:
    pct = score / max_score if max_score else 0
    if pct >= 0.70:
        return GREEN
    if pct >= 0.45:
        return AMBER
    return _rgb(0x9E, 0x9E, 0x9E)


# ── Slide builders ─────────────────────────────────────────────────────────────

def _slide_title(prs: Presentation, bp: Dict, ts: Dict, date_str: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _solid(slide.background, NAVY)

    # Red accent bar left
    _add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, RED)

    # Main title
    _add_textbox(slide, "M&A Target Screening",
                 Inches(0.5), Inches(1.6), Inches(9), Inches(1.0),
                 bold=True, size=38, color=WHITE)

    buyer  = bp.get("buyer", "Strategic Buyer")
    sector = ts.get("sector", "")
    geo    = ts.get("geography", "")

    _add_textbox(slide, buyer,
                 Inches(0.5), Inches(2.7), Inches(9), Inches(0.7),
                 bold=False, size=24, color=_rgb(0xC9, 0xA8, 0x4C))

    _add_textbox(slide, f"{sector}  ·  {geo}",
                 Inches(0.5), Inches(3.4), Inches(9), Inches(0.5),
                 bold=False, size=14, color=_rgb(0xAA, 0xB3, 0xC7))

    _add_textbox(slide, f"Prepared {date_str}  ·  Strictly Confidential",
                 Inches(0.5), Inches(6.6), Inches(9), Inches(0.4),
                 bold=False, size=11, color=_rgb(0x5A, 0x6A, 0x8A))

    # Right stats panel
    targets = ts.get("targets", [])
    _add_rect(slide, Inches(10.2), Inches(1.4), Inches(2.8), Inches(4.6), _rgb(0x1A, 0x1F, 0x42))
    stats = [
        ("Companies Screened", str(len(targets))),
        ("Scoring Criteria", "8"),
        ("Top Score", f"{targets[0].get('total_score','?')}/{targets[0].get('max_score','?')}" if targets else "—"),
        ("Recommendation", "PROCEED"),
    ]
    for i, (label, val) in enumerate(stats):
        y = Inches(1.7) + i * Inches(1.1)
        _add_textbox(slide, val,   Inches(10.4), y,              Inches(2.4), Inches(0.6),
                     bold=True, size=22, color=_rgb(0xC9, 0xA8, 0x4C))
        _add_textbox(slide, label, Inches(10.4), y + Inches(0.5), Inches(2.4), Inches(0.4),
                     bold=False, size=10, color=_rgb(0xAA, 0xB3, 0xC7))


def _slide_buyer(prs: Presentation, bp: Dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.28), RED)
    _heading(slide, f"Buyer Profile — {bp.get('buyer','')}", top=Inches(0.28))

    # Acquisitions — use buyer_acquisitions (new field) or fall back to acquisitions
    acqs = bp.get("buyer_acquisitions", bp.get("acquisitions", []))[:6]
    _add_textbox(slide, "Recent Acquisitions", Inches(0.4), Inches(1.0),
                 Inches(5.5), Inches(0.35), bold=True, size=12, color=NAVY)

    for i, a in enumerate(acqs):
        y = Inches(1.4) + i * Inches(0.45)
        name = str(a.get("name", ""))[:22]
        year = a.get("year", "")
        # support both old field name and new
        cap  = a.get("rationale", a.get("capability_gap_filled", ""))
        cap_short = str(cap)[:55]
        _add_textbox(slide, f"• {name} ({year})",
                     Inches(0.4), y, Inches(5.6), Inches(0.25),
                     bold=True, size=10, color=_rgb(0x2D, 0x37, 0x48))
        _add_textbox(slide, f"  {cap_short}",
                     Inches(0.4), y + Inches(0.24), Inches(5.6), Inches(0.22),
                     size=8, color=GREY)

    # Scoring criteria — show C1–C4 only (buyer-specific), keep descriptions short
    criteria = bp.get("scoring_criteria", [])[:4]
    _add_textbox(slide, "Scoring Criteria (C1–C4)", Inches(7.0), Inches(1.0),
                 Inches(5.8), Inches(0.35), bold=True, size=12, color=NAVY)

    for i, c in enumerate(criteria):
        y = Inches(1.4) + i * Inches(0.78)
        _add_rect(slide, Inches(7.0), y, Inches(0.42), Inches(0.6), RED)
        _add_textbox(slide, c.get("id",""), Inches(7.0), y + Inches(0.1),
                     Inches(0.42), Inches(0.4), bold=True, size=11, color=WHITE, align=PP_ALIGN.CENTER)
        _add_textbox(slide, str(c.get("name",""))[:40],
                     Inches(7.55), y, Inches(5.2), Inches(0.28),
                     bold=True, size=10, color=NAVY)
        desc = str(c.get("description",""))
        # First sentence only, capped at 90 chars
        first_sent = desc.split(".")[0][:90]
        _add_textbox(slide, first_sent,
                     Inches(7.55), y + Inches(0.28), Inches(5.2), Inches(0.42),
                     size=8, color=GREY)

    # Strategic priorities
    prios = bp.get("strategic_priorities", [])[:4]
    _add_textbox(slide, "Strategic Priorities", Inches(0.4), Inches(4.3),
                 Inches(5.5), Inches(0.35), bold=True, size=12, color=NAVY)
    for i, p in enumerate(prios):
        y = Inches(4.7) + i * Inches(0.38)
        _add_textbox(slide, f"→  {str(p)[:70]}", Inches(0.4), y, Inches(5.6), Inches(0.36),
                     size=10, color=_rgb(0x2D, 0x37, 0x48))


def _slide_rankings(prs: Presentation, ts: Dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.28), RED)
    _heading(slide, "Target Rankings — All Companies Screened", top=Inches(0.28))

    targets = ts.get("targets", [])
    criteria = ts.get("criteria_detail", [])
    max_score = targets[0].get("max_score", 35) if targets else 35

    col_x   = [Inches(0.3), Inches(0.7), Inches(2.8), Inches(4.5), Inches(5.8),
                Inches(7.1), Inches(8.4), Inches(9.7), Inches(11.0), Inches(12.3)]
    headers = ["#", "Company", "Country", "Stage", "C1", "C2", "C3", "C4", "Score", "Rec"]
    col_w   = [Inches(0.35), Inches(2.0), Inches(1.6), Inches(1.2),
               Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.1)]

    # Header row
    row_h = Inches(0.35)
    hy    = Inches(0.85)
    for hdr, x, w in zip(headers, col_x, col_w):
        _add_rect(slide, x, hy, w - Inches(0.04), row_h, NAVY)
        _add_textbox(slide, hdr, x + Inches(0.04), hy + Inches(0.04),
                     w - Inches(0.08), row_h - Inches(0.08),
                     bold=True, size=9, color=WHITE)

    # Data rows
    for i, t in enumerate(targets[:10]):
        y      = hy + row_h + i * Inches(0.55)
        scores = t.get("scores", {})
        rank   = t.get("rank", i+1)
        total  = t.get("total_score", 0)
        rec    = t.get("recommendation", "MONITOR")
        bg     = LIGHT if rank <= 3 else WHITE

        # Shorten funding stage to fit column
        stage = str(t.get("funding_stage",""))
        stage = stage.replace("Series ","Ser.").replace("Pre-Seed","Pre-S")[:10]

        def _score_cell(key):
            raw = scores.get(key, scores.get(key.lower(), "—"))
            if isinstance(raw, dict):
                return str(raw.get("score", "—"))
            return str(raw)

        score_str = f"{total}/{max_score}"
        row_vals = [
            str(rank),
            str(t.get("name",""))[:22],
            str(t.get("country",""))[:10],
            stage,
            _score_cell("C1"),
            _score_cell("C2"),
            _score_cell("C3"),
            _score_cell("C4"),
            score_str,
            rec[:7],
        ]
        for val, x, w in zip(row_vals, col_x, col_w):
            _add_rect(slide, x, y, w - Inches(0.04), Inches(0.5), bg)
            txt_color = _score_color(total, max_score) if val == score_str else NAVY
            _add_textbox(slide, val, x + Inches(0.04), y + Inches(0.04),
                         w - Inches(0.1), Inches(0.42), size=9, color=txt_color,
                         bold=(val == score_str))

        # Green left accent for top 3
        if rank <= 3:
            _add_rect(slide, col_x[0], y, Inches(0.04), Inches(0.5), GREEN)


def _slide_company(prs: Presentation, t: Dict, bp: Dict, rank: int,
                   detail: Optional[Dict] = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.28), RED)
    name = t.get("name", "")
    _heading(slide, f"#{rank} — {name}  ·  Acquisition Report", top=Inches(0.28))

    total     = t.get("total_score", 0)
    max_score = t.get("max_score", 35)
    scores    = t.get("scores", {})
    pct       = total / max_score if max_score else 0
    rec       = (detail or {}).get("recommendation", t.get("recommendation","MONITOR"))

    # Score badge
    score_color = _score_color(total, max_score)
    _add_rect(slide, Inches(10.8), Inches(0.9), Inches(2.2), Inches(1.1), score_color)
    _add_textbox(slide, f"{total}/{max_score}", Inches(10.8), Inches(0.95),
                 Inches(2.2), Inches(0.6), bold=True, size=28, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, "Strategic Fit Score", Inches(10.8), Inches(1.5),
                 Inches(2.2), Inches(0.35), bold=False, size=9, color=WHITE, align=PP_ALIGN.CENTER)

    # Recommendation badge
    rec_color = GREEN if "PROCEED" in rec.upper() else (AMBER if "MONITOR" in rec.upper() else RED)
    _add_rect(slide, Inches(10.8), Inches(2.1), Inches(2.2), Inches(0.5), rec_color)
    _add_textbox(slide, rec.upper()[:10], Inches(10.8), Inches(2.15),
                 Inches(2.2), Inches(0.4), bold=True, size=12, color=WHITE, align=PP_ALIGN.CENTER)

    # Company snapshot (left column)
    investors_raw = t.get("key_investors", [])
    if isinstance(investors_raw, list):
        investors_str = ", ".join(str(x) for x in investors_raw[:3])
    else:
        investors_str = str(investors_raw)
    investors_str = investors_str[:40]

    stage_raw = str(t.get("funding_stage","N/A"))
    stage_short = stage_raw[:20]

    snap = [
        ("Country",   str(t.get("country","N/A"))[:18]),
        ("Founded",   str(t.get("founded","N/A"))[:8]),
        ("Stage",     stage_short),
        ("Raised",    f"€{t.get('total_raised_usd_m','N/A')}M"),
        ("ARR",       f"€{t.get('arr_usd_m','N/A')}M"),
        ("Employees", str(t.get("employees","N/A"))[:10]),
        ("Investors", investors_str),
    ]
    _add_textbox(slide, "Company Snapshot", Inches(0.3), Inches(1.0),
                 Inches(4.0), Inches(0.35), bold=True, size=12, color=NAVY)
    for i, (label, val) in enumerate(snap):
        y = Inches(1.4) + i * Inches(0.42)
        _add_textbox(slide, label, Inches(0.3), y, Inches(1.5), Inches(0.38),
                     bold=True, size=9, color=GREY)
        _add_textbox(slide, str(val), Inches(1.85), y, Inches(2.4), Inches(0.38),
                     size=9, color=NAVY)

    # Criterion scores (middle column) — C1 through C8
    criteria = t.get("criteria_detail", bp.get("scoring_criteria", []))
    _add_textbox(slide, "Criterion Scores", Inches(4.6), Inches(1.0),
                 Inches(4.0), Inches(0.35), bold=True, size=12, color=NAVY)

    score_keys = [f"C{i}" for i in range(1, 9)]
    row_h_score = Inches(0.37)
    for i, ck in enumerate(score_keys):
        raw = scores.get(ck, scores.get(ck.lower()))
        if raw is None:
            continue
        score_val = raw.get("score", 1) if isinstance(raw, dict) else raw
        try:
            score_val = max(1, min(5, int(score_val)))
        except (TypeError, ValueError):
            score_val = 1
        y = Inches(1.4) + i * row_h_score
        cname = ""
        for c in criteria:
            if c.get("id","").upper() == ck:
                cname = c.get("name","")
                break
        label = (cname[:24] if cname else ck)
        _add_textbox(slide, label, Inches(4.6), y, Inches(3.4), Inches(row_h_score),
                     size=9, color=NAVY)
        bar_w = Inches(1.3) * score_val / 5
        _add_rect(slide, Inches(8.1), y + Inches(0.06), Inches(1.3), Inches(0.22), LIGHT)
        _add_rect(slide, Inches(8.1), y + Inches(0.06), bar_w, Inches(0.22),
                  _score_color(score_val, 5))
        _add_textbox(slide, str(score_val), Inches(9.5), y, Inches(0.38), Inches(row_h_score),
                     bold=True, size=9, color=NAVY)

    # Fit summary — two sentences max, hard cap at 220 chars
    fit_summary = ""
    if detail:
        fit_summary = detail.get("exec_summary", "")
        if not fit_summary:
            fit_summary = detail.get("advisory_view", "")
    if not fit_summary:
        fit_summary = t.get("strategic_fit_summary", "")
    # Keep first two sentences
    sentences = fit_summary.replace("  ", " ").split(". ")
    fit_short = ". ".join(sentences[:2])
    if len(fit_short) > 220:
        fit_short = fit_short[:217] + "..."

    _add_textbox(slide, "Strategic Fit Summary", Inches(0.3), Inches(4.75),
                 Inches(10.0), Inches(0.3), bold=True, size=11, color=NAVY)
    _add_rect(slide, Inches(0.3), Inches(5.1), Inches(10.0), Inches(1.5), _rgb(0xF8, 0xF9, 0xFF))
    _add_textbox(slide, fit_short if fit_short else "N/A",
                 Inches(0.45), Inches(5.15), Inches(9.7), Inches(1.4),
                 size=9, color=_rgb(0x2D, 0x37, 0x48), wrap=True)

    # Risks
    risks = t.get("deal_breaker_risks", [])
    if risks:
        _add_textbox(slide, "Key Risks", Inches(10.8), Inches(3.0),
                     Inches(2.2), Inches(0.35), bold=True, size=11, color=NAVY)
        for j, r in enumerate(risks[:3]):
            y = Inches(3.4) + j * Inches(0.5)
            _add_textbox(slide, f"⚠ {str(r)[:38]}", Inches(10.8), y,
                         Inches(2.2), Inches(0.46), size=8, color=_rgb(0xB7, 0x27, 0x27), wrap=True)


def _slide_dcf(prs: Presentation, dcf: Dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.28), RED)
    _heading(slide, f"DCF Valuation — {dcf.get('target_name','Top Target')}", top=Inches(0.28))

    # Assumptions panel
    assumptions = dcf.get("assumptions", {})
    _add_textbox(slide, "Model Assumptions", Inches(0.3), Inches(1.0),
                 Inches(3.5), Inches(0.35), bold=True, size=12, color=NAVY)
    _add_rect(slide, Inches(0.3), Inches(1.4), Inches(3.5), Inches(2.8), LIGHT)
    assum_rows = [
        ("Base Revenue",    f"€{assumptions.get('base_revenue_eur_m',700):.0f}M"),
        ("EBITDA Margin",   f"{assumptions.get('ebitda_margin_pct',15)}%"),
        ("FCF Conversion",  f"{assumptions.get('fcf_pct_ebitda',55)}%"),
        ("WACC",            f"{assumptions.get('wacc_pct',10.5)}%"),
        ("Terminal Growth", f"{assumptions.get('terminal_growth_pct',3.0)}%"),
        ("Synergy Premium", f"{assumptions.get('synergy_premium_pct',20)}%"),
    ]
    for i, (label, val) in enumerate(assum_rows):
        y = Inches(1.5) + i * Inches(0.42)
        _add_textbox(slide, label, Inches(0.4), y, Inches(1.9), Inches(0.36),
                     size=10, color=GREY)
        _add_textbox(slide, val, Inches(2.3), y, Inches(1.4), Inches(0.36),
                     bold=True, size=10, color=NAVY)

    # Valuation bridge
    bridge = [
        ("PV of FCFs",        dcf.get("sum_pv_fcf", 0),          NAVY),
        ("PV Terminal Value",  dcf.get("pv_terminal_value", 0),   _rgb(0x12,0x0A,0x8F)),
        ("Enterprise Value",   dcf.get("enterprise_value", 0),    GREEN),
        ("Synergy Premium",    dcf.get("synergy_value", 0),       AMBER),
        ("Acquisition Price",  dcf.get("acquisition_price", 0),   RED),
    ]
    max_val = max(v for _, v, _ in bridge) * 1.1 or 1
    _add_textbox(slide, "Valuation Bridge", Inches(4.2), Inches(1.0),
                 Inches(8.7), Inches(0.35), bold=True, size=12, color=NAVY)

    for i, (label, val, color) in enumerate(bridge):
        y    = Inches(1.4) + i * Inches(0.82)
        barw = Inches(7.0) * val / max_val
        _add_textbox(slide, label, Inches(4.2), y, Inches(2.5), Inches(0.35),
                     bold=True, size=11, color=NAVY)
        _add_rect(slide, Inches(6.8), y + Inches(0.04), Inches(7.0), Inches(0.42), LIGHT)
        _add_rect(slide, Inches(6.8), y + Inches(0.04), barw, Inches(0.42), color)
        _add_textbox(slide, f"€{val:,.0f}M", Inches(6.8) + barw + Inches(0.1), y,
                     Inches(2.0), Inches(0.4), bold=True, size=11, color=color)


def _slide_recommendation(prs: Presentation, ts: Dict, bp: Dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid(slide.background, NAVY)
    _add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, RED)

    _add_textbox(slide, "Recommendation & Next Steps",
                 Inches(0.5), Inches(0.6), Inches(12.0), Inches(0.7),
                 bold=True, size=26, color=WHITE)

    targets = ts.get("targets", [])
    top3    = [t for t in targets if t.get("rank",99) <= 3]

    # Top 3 recommendation cards
    for i, t in enumerate(top3):
        x = Inches(0.5) + i * Inches(4.2)
        _add_rect(slide, x, Inches(1.5), Inches(3.9), Inches(3.5), _rgb(0x1A,0x1F,0x42))
        _add_rect(slide, x, Inches(1.5), Inches(3.9), Inches(0.5), RED)
        _add_textbox(slide, f"#{t.get('rank','')}  {t.get('name','')}",
                     x + Inches(0.15), Inches(1.56), Inches(3.6), Inches(0.4),
                     bold=True, size=13, color=WHITE)
        _add_textbox(slide, t.get("country",""),
                     x + Inches(0.15), Inches(2.1), Inches(3.6), Inches(0.35),
                     size=10, color=_rgb(0xC9,0xA8,0x4C))
        score = f"{t.get('total_score','')}/{t.get('max_score','')}"
        _add_textbox(slide, f"Score: {score}",
                     x + Inches(0.15), Inches(2.5), Inches(3.6), Inches(0.35),
                     size=10, color=_rgb(0xAA,0xB3,0xC7))
        fit_raw = t.get("strategic_fit_summary","")
        fit_sents = fit_raw.replace("  ", " ").split(". ")
        fit = ". ".join(fit_sents[:2])
        if len(fit) > 160:
            fit = fit[:157] + "..."
        _add_textbox(slide, fit,
                     x + Inches(0.15), Inches(2.9), Inches(3.6), Inches(1.8),
                     size=8, color=_rgb(0xAA,0xB3,0xC7), wrap=True)

    # Next steps
    steps = [
        "1.  Conduct preliminary outreach to top-ranked targets via M&A advisor",
        "2.  Commission independent financial due diligence on #1 and #2 targets",
        "3.  Validate ARR, churn rate, and gross margin with management presentations",
        "4.  Engage legal counsel for NDA and exclusivity agreement drafting",
        "5.  Present shortlist to investment committee for formal mandate approval",
    ]
    _add_textbox(slide, "Recommended Next Steps",
                 Inches(0.5), Inches(5.2), Inches(12.5), Inches(0.4),
                 bold=True, size=14, color=_rgb(0xC9,0xA8,0x4C))
    for i, s in enumerate(steps):
        y = Inches(5.7) + i * Inches(0.32)
        _add_textbox(slide, s, Inches(0.5), y, Inches(12.5), Inches(0.3),
                     size=10, color=_rgb(0xAA,0xB3,0xC7))


# ── Public API ─────────────────────────────────────────────────────────────────

def build_pptx(bp: Dict, ts: Dict, dcf: Optional[Dict] = None,
               company_details: Optional[Dict] = None,
               output_dir: Optional[Path] = None) -> str:
    """Generate the deck and return the output path as a string."""
    _out_dir = Path(output_dir) if output_dir is not None else OUTPUT_DIR
    _out_dir.mkdir(parents=True, exist_ok=True)
    out_path = _out_dir / "report.pptx"

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    targets  = ts.get("targets", [])
    top3     = sorted([t for t in targets if t.get("rank", 99) <= 3],
                      key=lambda x: x.get("rank", 99))
    date_str = datetime.now().strftime("%d %B %Y")

    print("  [PPTX] Slide 1 — Title")
    _slide_title(prs, bp, ts, date_str)

    print("  [PPTX] Slide 2 — Buyer Profile")
    _slide_buyer(prs, bp)

    print("  [PPTX] Slide 3 — Rankings Table")
    _slide_rankings(prs, ts)

    for i, t in enumerate(top3, 1):
        print(f"  [PPTX] Slide {3+i} — {t.get('name','')}")
        detail = (company_details or {}).get(t.get("name", ""))
        _slide_company(prs, t, bp, rank=t.get("rank", i), detail=detail)

    print("  [PPTX] Slide — Recommendation")
    _slide_recommendation(prs, ts, bp)

    prs.save(str(out_path))
    print(f"  [PPTX] Saved → {out_path}")
    return str(out_path)


def run(bp: Optional[Dict] = None, ts: Optional[Dict] = None) -> str:
    if bp is None:
        with open(DATA_DIR / "buyer_profile.json", encoding="utf-8") as f:
            bp = json.load(f)
    if ts is None:
        with open(DATA_DIR / "targets_scored.json", encoding="utf-8") as f:
            ts = json.load(f)
    return build_pptx(bp, ts)


if __name__ == "__main__":
    print("\n[Step 5] PowerPoint Generation")
    print(run())
