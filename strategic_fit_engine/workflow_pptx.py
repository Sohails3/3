"""
Generates an editable PowerPoint version of the 3-slide workflow diagrams.

Slides:
  1 — M&A Screening Workflow (5-phase process)
  2 — Backend Architecture (API/tech stack)
  3 — Impact Analysis (comparison table)

Output: strategic_fit_engine/output/workflow_diagrams.pptx
"""
from __future__ import annotations

from pathlib import Path
from typing import Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Colours
# ---------------------------------------------------------------------------
NAVY   = RGBColor(0x25, 0x28, 0x50)
RED    = RGBColor(0xCC, 0x06, 0x05)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0xF1, 0xF5, 0xF9)
GREEN  = RGBColor(0x15, 0x80, 0x3D)
AMBER  = RGBColor(0x92, 0x40, 0x0E)
BLUE   = RGBColor(0x1D, 0x4E, 0xD8)
LIGHT  = RGBColor(0xE8, 0xEA, 0xF0)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
TEXT   = RGBColor(0x1E, 0x29, 0x3B)
GBORDER= RGBColor(0xD8, 0xDC, 0xE8)

OUTPUT = Path(__file__).parent / "output" / "workflow_diagrams.pptx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _inches(x: float) -> Emu:
    return Inches(x)


def _add_slide(prs: Presentation):
    blank = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank)


def _box(slide, l, t, w, h, fill=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(1, _inches(l), _inches(t), _inches(w), _inches(h))
    shape.line.width = line_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _text_box(slide, l, t, w, h, text, size=11, bold=False, color=TEXT,
              align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(_inches(l), _inches(t), _inches(w), _inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def _header(slide, tag: str, title: str, subtitle: str):
    """Navy header bar across top."""
    _box(slide, 0, 0, 13.33, 0.75, fill=NAVY)
    # Red top border
    _box(slide, 0, 0, 13.33, 0.045, fill=RED)
    # Tag pill
    _box(slide, 0.3, 0.17, 0.7, 0.22, fill=RED)
    _text_box(slide, 0.31, 0.175, 0.68, 0.2, tag.upper(),
              size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _text_box(slide, 1.1, 0.12, 9, 0.28, title,
              size=14, bold=True, color=WHITE)
    _text_box(slide, 1.1, 0.42, 9, 0.22, subtitle,
              size=9, color=RGBColor(0xAA, 0xAD, 0xC4))
    _text_box(slide, 11.5, 0.25, 1.6, 0.25,
              "GP Bullhound · Intelligence Engine",
              size=7, color=RGBColor(0x88, 0x8B, 0xAA), align=PP_ALIGN.RIGHT)


def _footer(slide, left_text: str):
    _box(slide, 0, 7.25, 13.33, 0.25, fill=NAVY)
    _text_box(slide, 0.3, 7.27, 7, 0.2, left_text,
              size=7, color=RGBColor(0x88, 0x8B, 0xAA))
    _text_box(slide, 9, 7.27, 4, 0.2, "Strictly Confidential  ·  April 2026",
              size=7, color=RGBColor(0x88, 0x8B, 0xAA), align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Slide 1 — Workflow
# ---------------------------------------------------------------------------

def _slide1(prs: Presentation):
    slide = _add_slide(prs)
    _header(slide, "Workflow", "M&A Target Screening — 5-Phase Process",
            "Buyer-first methodology · Criteria derived from competitive intelligence")

    phases = [
        ("01", "BUY-SIDE\nINTELLIGENCE", NAVY, [
            "M&A Strategy — stated priorities & thesis",
            "Dry Powder — cash, FCF, deal size range",
            "Acquisition History — last 5 deals mapped",
            "Competitor M&A — threats in target sector",
            "Market Signals — trends driving urgency",
        ]),
        ("02", "TARGET\nDISCOVERY", RGBColor(0x1A, 0x56, 0x8C), [
            "8 – 12 private companies identified",
            "Series A to pre-IPO stage",
            "UK / Germany / France / Nordics",
            "Funding, ARR, headcount extracted",
            "Readiness signals flagged per target",
        ]),
        ("03", "SCORING\nRUBRIC", RGBColor(0x1B, 0x6E, 0x4A), [
            "C1 – C4: Buyer-specific criteria",
            "C5: Technology & IP quality",
            "C6: Market position & moat",
            "C7: Team & talent retention risk",
            "C8: Legal & regulatory exposure",
        ]),
        ("04", "SHORTLIST\n& RISKS", RGBColor(0x7B, 0x34, 0x1E), [
            "Ranked by total score out of 40",
            "Top 3 / Mid tier / Deprioritise",
            "Deal-breaker risks flagged",
            "Strip profile — side-by-side comparison",
            "Acquisition readiness score 1 – 10",
        ]),
        ("05", "CLIENT\nOUTPUT", RGBColor(0x4C, 0x1D, 0x95), [
            "HTML report — browser & print",
            "Editable PowerPoint deck",
            "Indicative valuation ranges",
            "Data sources disclaimer",
            "Banker validates & presents",
        ]),
    ]

    x = 0.18
    w = 2.52
    gap = 0.08
    for i, (num, title, color, bullets) in enumerate(phases):
        # Phase card header
        _box(slide, x, 0.85, w, 0.58, fill=color)
        _text_box(slide, x + 0.1, 0.87, w - 0.2, 0.18,
                  f"PHASE {num}", size=7, bold=True,
                  color=RGBColor(0xAA, 0xAD, 0xC4))
        _text_box(slide, x + 0.1, 1.04, w - 0.2, 0.35,
                  title, size=11, bold=True, color=WHITE)
        # Phase card body
        _box(slide, x, 1.43, w, 5.7, fill=RGBColor(0xFA, 0xFB, 0xFF),
             line_color=GBORDER)
        for j, b in enumerate(bullets):
            _text_box(slide, x + 0.18, 1.52 + j * 0.44, w - 0.28, 0.4,
                      f"• {b}", size=9.5, color=TEXT)
        # Arrow between phases
        if i < len(phases) - 1:
            _text_box(slide, x + w + 0.005, 3.1, gap + 0.06, 0.3,
                      "▶", size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)
        x += w + gap + 0.06

    _footer(slide, "GP Bullhound  ·  Intelligence Engine — Screening Workflow")


# ---------------------------------------------------------------------------
# Slide 2 — Architecture
# ---------------------------------------------------------------------------

def _slide2(prs: Presentation):
    slide = _add_slide(prs)
    _header(slide, "Architecture", "Backend System — API & Technology Stack",
            "Claude API · FactSet · Flask · python-pptx")

    layers = [
        ("UI LAYER", "Flask Web App — Browser Interface",
         ["Buyer / Sector / Geography form inputs",
          "Real-time SSE streaming progress log",
          "Download: HTML report + PPTX deck"]),
        ("AI LAYER", "Anthropic Claude API (claude-sonnet-4-6)",
         ["Step 1: Web search — buyer financials, acquisitions, competitors",
          "Step 2: Training data — target discovery (10 companies)",
          "Step 3: Scoring — 8 criteria, batches of 5"]),
        ("DATA LAYER", "External Data Sources",
         ["FactSet API — verified financials, M&A transactions",
          "Companies House — UK incorporation & filing data",
          "Crunchbase / press — funding rounds & news"]),
        ("OUTPUT LAYER", "Report Generation",
         ["step4_output.py — HTML report (inline CSS, self-contained)",
          "step5_pptx.py — Editable PowerPoint via python-pptx",
          "output/archive/ — timestamped run history"]),
    ]

    colors = [NAVY, RGBColor(0x1A, 0x56, 0x8C), RGBColor(0x1B, 0x6E, 0x4A),
              RGBColor(0x4C, 0x1D, 0x95)]

    y = 0.9
    h_head = 0.32
    h_body = 0.95
    h_gap  = 0.18

    for i, (label, title, bullets) in enumerate(layers):
        c = colors[i]
        _box(slide, 0.3, y, 12.73, h_head, fill=c)
        _text_box(slide, 0.45, y + 0.05, 2, 0.22, label,
                  size=8, bold=True, color=RGBColor(0xAA, 0xAD, 0xC4))
        _text_box(slide, 2.2, y + 0.05, 10.6, 0.22, title,
                  size=11, bold=True, color=WHITE)
        _box(slide, 0.3, y + h_head, 12.73, h_body,
             fill=RGBColor(0xFA, 0xFB, 0xFF), line_color=GBORDER)
        for j, b in enumerate(bullets):
            _text_box(slide, 0.6, y + h_head + 0.08 + j * 0.28, 12.3, 0.26,
                      f"• {b}", size=10, color=TEXT)
        if i < len(layers) - 1:
            _text_box(slide, 6.5, y + h_head + h_body + 0.01, 0.5, 0.15,
                      "▼", size=11, bold=True, color=RED, align=PP_ALIGN.CENTER)
        y += h_head + h_body + h_gap

    _footer(slide, "GP Bullhound  ·  Intelligence Engine — Backend Architecture")


# ---------------------------------------------------------------------------
# Slide 3 — Impact Table
# ---------------------------------------------------------------------------

def _slide3(prs: Presentation):
    slide = _add_slide(prs)
    _header(slide, "Impact", "Value to Investment Banks & M&A Advisory",
            "Time saved · Cost reduction · Competitive advantage")

    # Summary metrics row
    metrics = [
        ("3 – 5 min", "Full pipeline runtime", NAVY, WHITE),
        ("~95%", "Reduction in longlist time", RGBColor(0xDC, 0xFC, 0xE7), GREEN),
        ("£8 – 15k", "Analyst hours saved per mandate", RGBColor(0xFE, 0xF3, 0xC7), AMBER),
        ("10×", "More mandates screened per analyst", RGBColor(0xEF, 0xF6, 0xFF), BLUE),
    ]
    mx = 0.3
    mw = 3.08
    for label, sub, bg, fg in metrics:
        _box(slide, mx, 0.9, mw, 0.72, fill=bg,
             line_color=GBORDER if bg != NAVY else None)
        _text_box(slide, mx + 0.1, 0.94, mw - 0.2, 0.38,
                  label, size=22, bold=True, color=WHITE if bg == NAVY else fg,
                  align=PP_ALIGN.CENTER)
        _text_box(slide, mx + 0.1, 1.3, mw - 0.2, 0.26,
                  sub, size=8, color=RGBColor(0xAA, 0xAD, 0xC4) if bg == NAVY else fg,
                  align=PP_ALIGN.CENTER)
        mx += mw + 0.09

    # Table header
    cols = [2.6, 2.3, 1.9, 1.6, 2.63]  # widths
    headers = ["Stage", "Traditional IB Process", "This Engine", "Time Saved", "Why It Matters"]
    hx = 0.3
    _box(slide, 0.3, 1.72, 12.73, 0.3, fill=NAVY)
    for w_col, h_txt in zip(cols, headers):
        _text_box(slide, hx + 0.06, 1.76, w_col - 0.1, 0.24,
                  h_txt.upper(), size=7.5, bold=True,
                  color=RGBColor(0xAA, 0xAD, 0xC4))
        hx += w_col

    rows = [
        ("Mandate Definition",
         "2–4 weeks stakeholder alignment",
         "2 minutes",
         "~3 weeks",
         "Buyer DNA derived automatically — no criteria debate"),
        ("Longlist Generation",
         "30–80 companies, 2–3 analyst days",
         "10–12 companies, ~90 sec",
         "2–3 days",
         "Frees senior analysts for relationships & execution"),
        ("Screening & Scoring",
         "PitchBook/CapIQ pull + manual scoring, 1–2 weeks",
         "8-criteria rubric, ~60 sec",
         "1–2 weeks",
         "Consistent, bias-free scoring across all targets"),
        ("Competitor Mapping",
         "Senior banker research, 3–5 days",
         "Automated in Step 1, ~60 sec",
         "3–5 days",
         "Threat-driven criteria reflect real deal urgency"),
        ("Client-Ready Output",
         "PowerPoint deck, 1–2 days formatting",
         "HTML + PPTX auto-generated",
         "1–2 days",
         "Faster to client — competitive advantage in pitches"),
        ("Off-Market Sourcing",
         "Banker network — major advantage",
         "Not covered",
         "No saving",
         "Relationship-driven sourcing remains the banker's edge"),
    ]

    row_colors = [RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0xF9, 0xFA, 0xFB)]
    amber_row = RGBColor(0xFE, 0xF9, 0xF0)

    ry = 2.02
    rh = 0.74
    for ri, (stage, trad, engine, saved, why) in enumerate(rows):
        bg = amber_row if ri == 5 else row_colors[ri % 2]
        _box(slide, 0.3, ry, 12.73, rh, fill=bg, line_color=GBORDER)
        rx = 0.3
        cell_data = [stage, trad, engine, saved, why]
        for ci, (w_col, cell) in enumerate(zip(cols, cell_data)):
            is_engine = ci == 2
            is_saved  = ci == 3
            is_stage  = ci == 0
            color = GREEN if (is_engine and ri < 5) else \
                    AMBER  if (is_engine and ri == 5) else \
                    GREEN  if (is_saved and ri < 5) else \
                    AMBER  if (is_saved and ri == 5) else \
                    NAVY   if is_stage else TEXT
            bold = is_stage or is_engine
            _text_box(slide, rx + 0.06, ry + 0.06, w_col - 0.1, rh - 0.1,
                      cell, size=8.5, bold=bold, color=color)
            rx += w_col
        ry += rh

    # Bottom note
    _box(slide, 0.3, 6.54, 12.73, 0.58, fill=GREY, line_color=GBORDER)
    _box(slide, 0.3, 6.54, 0.05, 0.58, fill=RED)
    _text_box(slide, 0.5, 6.6, 12.4, 0.46,
              "This engine augments bankers — it does not replace them. "
              "Relationship-driven off-market sourcing, sellability assessment, and negotiation "
              "remain exclusively human. The engine compresses 3–6 weeks of analytical groundwork "
              "into minutes, letting senior bankers focus on where they create the most value.",
              size=9, color=TEXT)

    _footer(slide, "GP Bullhound  ·  Intelligence Engine — Impact Analysis")


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def _slide4(prs: Presentation):
    slide = _add_slide(prs)
    _header(slide, "System", "How the Strategic Fit Engine Works",
            "Tools · Data sources · Workflow steps")

    # Three column layout
    col_w   = 4.0
    col_gap = 0.2
    col_y   = 0.95
    col_h   = 6.1

    # ── Column 1: Tools Used ──
    cx = 0.3
    _box(slide, cx, col_y, col_w, 0.28, fill=NAVY)
    _text_box(slide, cx + 0.1, col_y + 0.05, col_w - 0.2, 0.2,
              "TOOLS USED", size=8, bold=True, color=RGBColor(0xAA, 0xAD, 0xC4))

    tools = [
        ("Anthropic Claude API",
         "Model: claude-sonnet-4-6. Handles buyer research (web search), "
         "target discovery, scoring, and report generation."),
        ("Flask (Python)",
         "Browser-based UI with real-time progress streaming via "
         "Server-Sent Events. Runs locally on port 8080."),
        ("python-pptx",
         "Generates fully editable PowerPoint decks programmatically — "
         "no manual formatting required."),
        ("pdfplumber / openpyxl",
         "Optional: extracts verified financials from uploaded PDF annual "
         "reports or Excel spreadsheets to anchor model outputs."),
    ]

    ty = col_y + 0.32
    for name, desc in tools:
        _box(slide, cx, ty, col_w, 1.16, fill=RGBColor(0xF8, 0xF9, 0xFF),
             line_color=GBORDER)
        _text_box(slide, cx + 0.1, ty + 0.08, col_w - 0.2, 0.22,
                  name, size=10, bold=True, color=NAVY)
        _text_box(slide, cx + 0.1, ty + 0.3, col_w - 0.2, 0.8,
                  desc, size=9, color=TEXT)
        ty += 1.2

    # ── Column 2: Data Sources ──
    cx = 0.3 + col_w + col_gap
    _box(slide, cx, col_y, col_w, 0.28, fill=RED)
    _text_box(slide, cx + 0.1, col_y + 0.05, col_w - 0.2, 0.2,
              "DATA SOURCES", size=8, bold=True, color=RGBColor(0xFF, 0xCC, 0xCC))

    sources = [
        ("Live — Web Search", RGBColor(0xFF, 0xF8, 0xF8), RGBColor(0x99, 0x1B, 0x1B),
         ["Company websites & press releases",
          "TechCrunch, Sifted, EU-Startups",
          "Crunchbase public profiles",
          "Companies House (UK filings)",
          "Buyer investor letters & earnings"]),
        ("Verified — Optional Upload", RGBColor(0xF0, 0xFD, 0xF4), GREEN,
         ["Buyer annual report / 10-K PDF",
          "Excel spreadsheet — ARR & raised",
          "FactSet API (if provisioned)",
          "Capital IQ Web Services (if provisioned)"]),
        ("AI-Synthesised — Training Data", RGBColor(0xFE, 0xFC, 0xE8), AMBER,
         ["Target company profiles (Step 2)",
          "Headcount estimates from LinkedIn",
          "ARR estimated from stage + headcount",
          "Flagged in report — verify independently"]),
    ]

    ty = col_y + 0.32
    for label, bg, fg, bullets in sources:
        h = 0.28 + len(bullets) * 0.26
        _box(slide, cx, ty, col_w, h, fill=bg, line_color=GBORDER)
        _text_box(slide, cx + 0.1, ty + 0.06, col_w - 0.2, 0.2,
                  label, size=9.5, bold=True, color=fg)
        for bi, b in enumerate(bullets):
            _text_box(slide, cx + 0.1, ty + 0.28 + bi * 0.26, col_w - 0.2, 0.24,
                      f"• {b}", size=9, color=TEXT)
        ty += h + 0.1

    # ── Column 3: Workflow Steps ──
    cx = 0.3 + (col_w + col_gap) * 2
    _box(slide, cx, col_y, col_w, 0.28, fill=BLUE)
    _text_box(slide, cx + 0.1, col_y + 0.05, col_w - 0.2, 0.2,
              "WORKFLOW STEPS", size=8, bold=True, color=RGBColor(0xBF, 0xD7, 0xFF))

    steps = [
        ("1", "Buyer Intelligence", "~60s",
         "Web search fetches live financials, last 5 acquisitions, competitor "
         "M&A. Derives 4 buyer-specific scoring criteria (C1–C4)."),
        ("2", "Target Discovery", "~90s",
         "Identifies 10–12 private companies (Series A–pre-IPO) in the target "
         "sector and geography. Extracts funding, ARR, headcount, readiness signals."),
        ("3", "Scoring", "~60s",
         "Scores each company against 8 criteria (C1–C4 buyer-specific + "
         "C5–C8 universal). Python recalculates all totals. Ranked out of 40."),
        ("4", "Report Generation", "~30s",
         "Self-contained HTML report + editable PPTX. Ranked table, top 3 "
         "recommendation cards, deal-breaker risks, data sources disclaimer."),
    ]

    ty = col_y + 0.32
    for num, title, timing, desc in steps:
        _box(slide, cx, ty, col_w, 1.3, fill=RGBColor(0xF8, 0xF9, 0xFF),
             line_color=GBORDER)
        # Step number badge
        _box(slide, cx + 0.1, ty + 0.14, 0.3, 0.3, fill=NAVY)
        _text_box(slide, cx + 0.1, ty + 0.14, 0.3, 0.3,
                  num, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _text_box(slide, cx + 0.48, ty + 0.08, col_w - 0.65, 0.22,
                  title, size=10.5, bold=True, color=NAVY)
        _text_box(slide, cx + 0.48, ty + 0.28, 1.5, 0.18,
                  timing, size=8.5, color=MUTED)
        _text_box(slide, cx + 0.1, ty + 0.5, col_w - 0.2, 0.72,
                  desc, size=9, color=TEXT)
        ty += 1.38

    # Inputs/Outputs bar
    _box(slide, cx, ty, col_w, 0.5, fill=GREY, line_color=GBORDER)
    _box(slide, cx, ty, 0.04, 0.5, fill=RED)
    _text_box(slide, cx + 0.12, ty + 0.05, col_w - 0.2, 0.18,
              "Inputs: Buyer name · Target sector · Geography", size=8.5,
              bold=True, color=NAVY)
    _text_box(slide, cx + 0.12, ty + 0.24, col_w - 0.2, 0.18,
              "Outputs: report.html · report.pptx · targets_scored.json", size=8.5,
              color=TEXT)

    _footer(slide, "GP Bullhound  ·  Intelligence Engine — System Overview")


def run() -> Path:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    _slide1(prs)
    _slide2(prs)
    _slide3(prs)
    _slide4(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"  [PPTX] Workflow diagrams saved → {OUTPUT}")
    return OUTPUT


if __name__ == "__main__":
    path = run()
    print(f"Done: {path}")
